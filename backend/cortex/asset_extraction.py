"""Asset text extraction — pluggable per-type dispatcher.

Each kind (pdf/image/url) returns a normalized dict:

  {
    "text":     "<full extracted text, capped at ~12k chars>",
    "meta":     {<kind-specific metadata: dimensions, page_count, etc.>},
    "thumb_b64": "<optional base64 PNG thumbnail for image previews>",
  }

This module deliberately does NO LLM work — that lives in
`cortex.asset_intelligence`. Keeping the extraction step pure makes
the pipeline testable and lets us add PPTX/Video later by registering
a new entry in `_EXTRACTORS`.
"""
from __future__ import annotations

import asyncio
import base64
import io
import logging
import re
from typing import Optional

logger = logging.getLogger(__name__)

# Maximum text passed downstream to the LLM — ~12k chars ≈ 3k tokens.
# Anything more would inflate the prompt without adding insight.
TEXT_CAP = 12_000


async def extract(*, kind: str, data: bytes | None = None,
                    url: str | None = None) -> dict:
    """Dispatch to the per-kind extractor. Always returns a dict; on
    failure returns {"text": "", "meta": {"error": <reason>}} so the
    pipeline can still produce a partial brief instead of crashing."""
    fn = _EXTRACTORS.get(kind)
    if not fn:
        return {"text": "", "meta": {"error": f"unknown_kind:{kind}"}}
    try:
        if kind == "url":
            return await fn(url or "")
        return await fn(data or b"")
    except Exception as e:
        logger.exception("asset_extraction: %s extractor crashed", kind)
        return {"text": "", "meta": {"error": f"{type(e).__name__}: {e}"}}


# --------------------------------------------------------------- PDF
async def _extract_pdf(data: bytes) -> dict:
    """Pull text + page count using PyMuPDF (offline, no API calls)."""
    def _run():
        import fitz  # PyMuPDF
        doc = fitz.open(stream=data, filetype="pdf")
        pages: list[str] = []
        for i in range(min(doc.page_count, 30)):  # cap at 30 pages
            try:
                pages.append(doc.load_page(i).get_text("text"))
            except Exception:
                continue
        text = "\n\n".join(p for p in pages if p.strip())
        return {
            "text": text[:TEXT_CAP],
            "meta": {
                "page_count":      doc.page_count,
                "pages_extracted": min(doc.page_count, 30),
                "char_count":      len(text),
            },
        }
    return await asyncio.to_thread(_run)


# -------------------------------------------------------------- Image
async def _extract_image(data: bytes) -> dict:
    """Pull dimensions + a small thumbnail. We deliberately skip OCR
    here — the LLM with vision capability ingests the raw image during
    intelligence extraction, so OCR adds latency without new signal."""
    def _run():
        from PIL import Image
        img = Image.open(io.BytesIO(data))
        img.load()
        w, h = img.size
        fmt = (img.format or "").lower()
        # Thumb: 320px wide, preserve aspect, PNG for transparency safety.
        thumb = img.copy()
        thumb.thumbnail((320, 320), Image.LANCZOS)
        buf = io.BytesIO()
        thumb.convert("RGB").save(buf, format="JPEG", quality=78,
                                    optimize=True)
        thumb_b64 = base64.b64encode(buf.getvalue()).decode("ascii")
        return {
            # No text for raw images — the LLM gets the original bytes.
            "text": "",
            "meta": {
                "width":  w, "height": h,
                "format": fmt or "jpeg",
                "thumb_format": "jpeg",
            },
            "thumb_b64": thumb_b64,
        }
    return await asyncio.to_thread(_run)


# ---------------------------------------------------------------- URL
async def _extract_url(url: str) -> dict:
    """Fetch a webpage and return a clean text excerpt + title.
    Re-uses the existing site-fetcher pattern from analysis_runner.
    Returns empty text on fetch failure (caller decides if that's fatal).
    """
    if not url or not re.match(r"^https?://", url):
        return {"text": "", "meta": {"error": "invalid_url"}}
    try:
        import httpx
        async with httpx.AsyncClient(follow_redirects=True,
                                        timeout=15.0,
                                        headers={"User-Agent":
                                                  "Mozilla/5.0 CortexViral/1.0"}) as client:
            r = await client.get(url)
            html = r.text or ""
    except Exception as e:
        return {"text": "", "meta": {"error": f"fetch_failed:{type(e).__name__}"}}

    # Simple HTML → text. We don't need beautifulsoup precision; the
    # LLM is robust to noise. Just strip script/style and collapse tags.
    cleaned = re.sub(r"(?is)<(script|style|svg|noscript)\b[^>]*>.*?</\1>",
                       " ", html)
    cleaned = re.sub(r"(?is)<!--.*?-->", " ", cleaned)
    text = re.sub(r"(?s)<[^>]+>", " ", cleaned)
    text = re.sub(r"\s+", " ", text).strip()

    # Pull a title for the asset name fallback.
    title_m = re.search(r"(?is)<title[^>]*>(.*?)</title>", html)
    title = re.sub(r"\s+", " ", (title_m.group(1) if title_m else "")).strip()

    return {
        "text": text[:TEXT_CAP],
        "meta": {
            "url":         url,
            "title":       title[:200],
            "html_bytes":  len(html),
            "text_chars":  len(text),
        },
    }


# --------------------------------------------------------------- PPTX
async def _extract_pptx(data: bytes) -> dict:
    """Pull slide titles, bullet text, speaker notes, and a 'hero
    slide' rendering from a PPTX deck (offline, no API calls).

    The hero slide picker scores every slide on (a) slide-1 bias,
    (b) total image area, (c) presence of a title. The winning slide
    is rendered to a base64 PNG that downstream consumers (Phase B
    image gen, Assets card thumbnail) can use as a visual brand
    reference."""
    def _run():
        from pptx import Presentation
        from pptx.util import Emu
        from PIL import Image, ImageDraw, ImageFont
        prs = Presentation(io.BytesIO(data))
        slides: list[str] = []
        total_shapes = 0
        slide_meta: list[dict] = []
        # Map slide index → list of (image_bytes, ext) tuples so we can
        # surface the largest embedded picture too.
        slide_images: dict[int, list[tuple[bytes, str]]] = {}

        for i, slide in enumerate(prs.slides):
            chunks: list[str] = []
            title = ""
            try:
                if slide.shapes.title and slide.shapes.title.has_text_frame:
                    title = (slide.shapes.title.text_frame.text or "").strip()
            except Exception:
                title = ""
            if title:
                chunks.append(f"# Slide {i+1}: {title}")
            else:
                chunks.append(f"# Slide {i+1}")

            image_area = 0
            for shape in slide.shapes:
                total_shapes += 1
                # Collect embedded pictures (shape_type=13 = PICTURE).
                if getattr(shape, "shape_type", None) == 13:
                    try:
                        blob = shape.image.blob
                        ext = (shape.image.ext or "png").lower()
                        slide_images.setdefault(i, []).append((blob, ext))
                        # rough rendered area in EMU² for hero-pick
                        w = int(getattr(shape, "width", 0) or 0)
                        h = int(getattr(shape, "height", 0) or 0)
                        image_area += max(w * h, 0)
                    except Exception:
                        pass
                if not getattr(shape, "has_text_frame", False):
                    continue
                txt = (shape.text_frame.text or "").strip()
                if not txt or txt == title:
                    continue
                chunks.append(txt)
            try:
                notes = slide.notes_slide.notes_text_frame.text \
                    if slide.has_notes_slide else ""
                notes = (notes or "").strip()
                if notes:
                    chunks.append(f"[Speaker notes] {notes}")
            except Exception:
                pass

            slide_meta.append({
                "index":       i,
                "has_title":   bool(title),
                "image_area":  image_area,
                "image_count": len(slide_images.get(i, [])),
                "text_chars":  sum(len(c) for c in chunks),
            })
            slides.append("\n".join(chunks))

        text = "\n\n".join(slides).strip()

        # Hero-slide pick: slide-1 bias (×3) + normalized image area + title bonus.
        # Result is a deterministic 'most visually-loaded' slide.
        max_area = max((m["image_area"] for m in slide_meta), default=1) or 1
        for m in slide_meta:
            m["hero_score"] = (
                (3.0 if m["index"] == 0 else 0.0)
                + (m["image_area"] / max_area) * 2.0
                + (0.5 if m["has_title"] else 0.0)
                + min(m["image_count"], 3) * 0.3
            )
        hero = max(slide_meta, key=lambda m: m["hero_score"], default=None) \
            if slide_meta else None

        # Render the hero slide: prefer the largest embedded image
        # (highest resolution brand reference). Fall back to a text-
        # rendered placeholder card so the asset card still has a thumb.
        hero_b64 = ""
        hero_format = "png"
        if hero is not None:
            imgs = slide_images.get(hero["index"]) or []
            if imgs:
                imgs.sort(key=lambda x: len(x[0]), reverse=True)
                blob, ext = imgs[0]
                try:
                    img = Image.open(io.BytesIO(blob))
                    img.load()
                    img.thumbnail((480, 480), Image.LANCZOS)
                    buf = io.BytesIO()
                    img.convert("RGB").save(buf, format="JPEG",
                                              quality=82, optimize=True)
                    hero_b64 = base64.b64encode(buf.getvalue()).decode("ascii")
                    hero_format = "jpeg"
                except Exception:
                    logger.exception("pptx hero image render failed")
            if not hero_b64:
                # Synthesize a card-style PNG from the slide title +
                # bullets so the Asset gallery still has a thumb.
                try:
                    card = Image.new("RGB", (480, 270), (16, 18, 24))
                    drw = ImageDraw.Draw(card)
                    raw = (slides[hero["index"]] or "").splitlines()
                    lines = [ln.lstrip("# ").strip() for ln in raw if ln.strip()][:5]
                    y = 30
                    for ln in lines:
                        drw.text((28, y), ln[:55], fill=(230, 232, 240))
                        y += 22
                    buf = io.BytesIO()
                    card.save(buf, format="PNG", optimize=True)
                    hero_b64 = base64.b64encode(buf.getvalue()).decode("ascii")
                    hero_format = "png"
                except Exception:
                    logger.exception("pptx hero placeholder render failed")

        return {
            "text": text[:TEXT_CAP],
            "meta": {
                "slide_count":   len(prs.slides),
                "shape_count":   total_shapes,
                "char_count":    len(text),
                "hero_slide_index":  (hero or {}).get("index"),
                "hero_slide_score":  round((hero or {}).get("hero_score", 0), 3),
                "hero_image_count":  sum(len(v) for v in slide_images.values()),
                "thumb_format":      hero_format if hero_b64 else None,
            },
            "thumb_b64": hero_b64,
        }
    return await asyncio.to_thread(_run)


# --------------------------------------------------------------- Video
async def _extract_video(data: bytes) -> dict:
    """Extract a first-second keyframe (→ thumb_b64) and the full
    transcript (via Whisper, Emergent LLM Key) from a short video clip.

    Caps:
      * input bytes ≤ 50 MiB (enforced at route layer too).
      * duration ≤ 5 min — anything longer gets the first 5 min only.
      * audio extracted at 32 kbps mono so the Whisper upload stays
        well under the 25 MiB API limit.
    """
    import os
    import shutil
    import subprocess
    import tempfile

    def _run() -> dict:
        ffmpeg_bin: Optional[str] = None
        try:
            import imageio_ffmpeg  # type: ignore
            ffmpeg_bin = imageio_ffmpeg.get_ffmpeg_exe()
        except Exception:
            ffmpeg_bin = shutil.which("ffmpeg")
        if not ffmpeg_bin:
            return {"text": "",
                    "meta":  {"error": "ffmpeg_unavailable"}}

        tmpdir = tempfile.mkdtemp(prefix="cortex_video_")
        try:
            src_path  = os.path.join(tmpdir, "src")
            thumb_png = os.path.join(tmpdir, "thumb.png")
            audio_mp3 = os.path.join(tmpdir, "audio.mp3")
            with open(src_path, "wb") as fh:
                fh.write(data)

            # 1) Single keyframe at t=1s (or t=0 for very short clips).
            #    Scale to width 320 preserving aspect (matches image asset
            #    thumb convention) so the Assets grid renders consistently.
            try:
                subprocess.run(
                    [ffmpeg_bin, "-y", "-loglevel", "error",
                     "-ss", "1", "-i", src_path,
                     "-vframes", "1", "-vf", "scale=320:-2",
                     thumb_png],
                    check=True, timeout=30)
            except Exception:
                # t=1s failed (very short clip?). Retry at t=0.
                try:
                    subprocess.run(
                        [ffmpeg_bin, "-y", "-loglevel", "error",
                         "-i", src_path,
                         "-vframes", "1", "-vf", "scale=320:-2",
                         thumb_png],
                        check=True, timeout=30)
                except Exception:
                    pass

            thumb_b64 = ""
            if os.path.exists(thumb_png):
                with open(thumb_png, "rb") as fh:
                    thumb_b64 = base64.b64encode(fh.read()).decode("ascii")

            # 2) Audio → mono 32 kbps mp3, capped at 5 min.
            try:
                subprocess.run(
                    [ffmpeg_bin, "-y", "-loglevel", "error",
                     "-i", src_path,
                     "-vn", "-ac", "1", "-ar", "16000",
                     "-b:a", "32k", "-t", "300",
                     audio_mp3],
                    check=True, timeout=180)
            except Exception:
                pass

            # 3) Probe duration for meta.
            duration_s = None
            try:
                probe = subprocess.run(
                    [ffmpeg_bin, "-i", src_path, "-hide_banner"],
                    capture_output=True, text=True, timeout=15)
                # ffmpeg emits to stderr; parse "Duration: HH:MM:SS.XX"
                import re as _re
                m = _re.search(r"Duration:\s+(\d+):(\d+):([\d.]+)",
                                  probe.stderr or "")
                if m:
                    h, mn, sec = int(m.group(1)), int(m.group(2)), float(m.group(3))
                    duration_s = h * 3600 + mn * 60 + sec
            except Exception:
                pass

            audio_path: Optional[str] = audio_mp3 if os.path.exists(audio_mp3) else None
            audio_size = os.path.getsize(audio_path) if audio_path else 0
            return {
                "thumb_b64":  thumb_b64,
                "audio_path": audio_path,
                "duration_s": duration_s,
                "audio_size": audio_size,
            }
        finally:
            # We DO need the audio file outside _run() for the async
            # Whisper call, so we postpone tmpdir cleanup to the caller.
            pass

    extracted = await asyncio.to_thread(_run)
    if extracted.get("meta", {}).get("error"):
        return extracted

    # Transcribe outside the thread so we don't block. Whisper accepts
    # mp3 directly per playbook; 25 MiB ceiling is way above our 32 k
    # mono * 5 min ≈ 1.2 MiB output.
    transcript = ""
    audio_path = extracted.get("audio_path")
    if audio_path:
        try:
            from emergentintegrations.llm.openai import OpenAISpeechToText
            import os as _os
            stt = OpenAISpeechToText(api_key=_os.environ["EMERGENT_LLM_KEY"])
            with open(audio_path, "rb") as fh:
                resp = await stt.transcribe(
                    file=fh, model="whisper-1",
                    response_format="json")
            transcript = (getattr(resp, "text", "") or "").strip()
        except Exception:
            logger.exception("video transcript failed")
        finally:
            # Drop the temp audio (+ keep the dir; OS reaper handles it).
            try:
                import os as _os2
                _os2.unlink(audio_path)
            except Exception:
                pass

    return {
        "text": transcript[:TEXT_CAP],
        "meta": {
            "duration_s":   extracted.get("duration_s"),
            "audio_bytes":  extracted.get("audio_size"),
            "char_count":   len(transcript),
            "thumb_format": "png" if extracted.get("thumb_b64") else None,
            # Whisper-1 list price is $0.006 / audio minute. We bill
            # the duration we transcribed (capped at 5 min), not the
            # original clip length, so the number matches reality.
            "transcription_cost_usd": (
                round(min(extracted.get("duration_s") or 0, 300) / 60 * 0.006, 4)
                if transcript else 0.0
            ),
            "transcription_provider": "whisper-1" if transcript else None,
        },
        "thumb_b64": extracted.get("thumb_b64") or "",
    }


_EXTRACTORS = {
    "pdf":   _extract_pdf,
    "image": _extract_image,
    "url":   _extract_url,
    "pptx":  _extract_pptx,
    "video": _extract_video,
}


def kind_from_mime(mime: str | None) -> Optional[str]:
    """Normalize an HTTP content-type into our asset kind enum."""
    if not mime:
        return None
    m = mime.lower()
    if m == "application/pdf":
        return "pdf"
    if m.startswith("image/"):
        return "image"
    if m == ("application/vnd.openxmlformats-officedocument"
             ".presentationml.presentation"):
        return "pptx"
    if m.startswith("video/"):
        return "video"
    return None
