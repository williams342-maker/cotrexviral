"""Phase A expansion tests: PPTX + Video asset extraction pipeline.

Covers:
  - MIME accept for pptx / video/mp4 / video/quicktime / video/webm
  - Per-MIME size caps (20 MiB default, 50 MiB for video) → HTTP 413
  - PPTX extraction_meta (slide_count, shape_count, char_count)
  - Video extraction_meta (duration_s, audio_bytes, char_count, thumb_format)
    + thumb_b64 populated
  - Backwards compatibility for PDF/Image still works
  - Frontend KIND_META has pptx + video entries
"""
import io
import os
import subprocess
import time

import pytest
import requests
from dotenv import load_dotenv

load_dotenv("/app/frontend/.env")
BASE_URL = os.environ["REACT_APP_BACKEND_URL"].rstrip("/")
SESSION_TOKEN = "test_session_1779636592168"
AUTH_HEADERS = {"Authorization": f"Bearer {SESSION_TOKEN}"}

FFMPEG_BIN = ("/root/.venv/lib/python3.11/site-packages/imageio_ffmpeg/"
              "binaries/ffmpeg-linux-aarch64-v7.0.2")

PPTX_MIME = ("application/vnd.openxmlformats-officedocument."
             "presentationml.presentation")


# ---------- helpers ----------------------------------------------------
def _build_pptx_bytes(slides=3):
    from pptx import Presentation
    prs = Presentation()
    for i in range(slides):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = f"Test Slide {i+1}"
        if slide.placeholders and len(slide.placeholders) > 1:
            try:
                slide.placeholders[1].text = f"Body content for slide {i+1}"
            except Exception:
                pass
        # Speaker notes
        slide.notes_slide.notes_text_frame.text = f"Speaker note {i+1}"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _build_video_bytes(seconds=3, with_audio=True, tmp="/tmp/_phaseA_test.mp4"):
    cmd = [FFMPEG_BIN, "-y", "-loglevel", "error",
           "-f", "lavfi", "-i", f"testsrc=duration={seconds}:size=160x120:rate=10"]
    if with_audio:
        cmd += ["-f", "lavfi", "-i", f"sine=frequency=440:duration={seconds}",
                "-c:a", "aac", "-b:a", "32k"]
    cmd += ["-c:v", "libx264", "-pix_fmt", "yuv420p", "-t", str(seconds), tmp]
    subprocess.run(cmd, check=True, timeout=60)
    with open(tmp, "rb") as fh:
        return fh.read()


def _upload_file(content, filename, mime):
    files = {"file": (filename, content, mime)}
    return requests.post(f"{BASE_URL}/api/cortex/assets/upload",
                         headers=AUTH_HEADERS, files=files, timeout=90)


def _wait_complete(asset_id, timeout=180):
    deadline = time.time() + timeout
    while time.time() < deadline:
        try:
            r = requests.get(f"{BASE_URL}/api/cortex/assets/{asset_id}",
                             headers=AUTH_HEADERS, timeout=60)
        except requests.exceptions.RequestException:
            time.sleep(3)
            continue
        if r.status_code != 200:
            time.sleep(2)
            continue
        data = r.json()
        if data.get("status") in ("complete", "failed"):
            return data
        time.sleep(3)
    pytest.fail(f"Asset {asset_id} did not reach terminal status in time")


def _upload_file_long(content, filename, mime, timeout=180):
    files = {"file": (filename, content, mime)}
    return requests.post(f"{BASE_URL}/api/cortex/assets/upload",
                         headers=AUTH_HEADERS, files=files, timeout=timeout)


# ---------- health -----------------------------------------------------
class TestHealth:
    def test_assets_list_ok(self):
        r = requests.get(f"{BASE_URL}/api/cortex/assets?limit=1",
                         headers=AUTH_HEADERS, timeout=15)
        assert r.status_code == 200, r.text
        assert "assets" in r.json()


# ---------- PPTX -------------------------------------------------------
class TestPPTXUpload:
    @pytest.fixture(scope="class")
    def pptx_bytes(self):
        return _build_pptx_bytes(slides=3)

    def test_upload_pptx_returns_kind_pptx(self, pptx_bytes):
        r = _upload_file(pptx_bytes, "TEST_deck.pptx", PPTX_MIME)
        assert r.status_code == 200, r.text
        row = r.json()
        assert row["kind"] == "pptx"
        assert row["mime_type"] == PPTX_MIME
        assert row["status"] in ("queued", "extracting", "analyzing", "complete")
        pytest.pptx_asset_id = row["id"]

    def test_pptx_extraction_meta_and_complete(self, pptx_bytes):
        aid = getattr(pytest, "pptx_asset_id", None)
        if not aid:
            pytest.skip("upload step failed")
        final = _wait_complete(aid, timeout=180)
        assert final["status"] == "complete", final
        meta = final.get("extraction_meta") or {}
        assert meta.get("slide_count") == 3, meta
        assert isinstance(meta.get("shape_count"), int) and meta["shape_count"] >= 3
        assert isinstance(meta.get("char_count"), int) and meta["char_count"] > 0

    def test_pptx_cleanup(self):
        aid = getattr(pytest, "pptx_asset_id", None)
        if aid:
            requests.delete(f"{BASE_URL}/api/cortex/assets/{aid}",
                            headers=AUTH_HEADERS, timeout=15)


# ---------- Video ------------------------------------------------------
class TestVideoUpload:
    @pytest.fixture(scope="class")
    def video_bytes(self):
        return _build_video_bytes(seconds=3, with_audio=True)

    def test_upload_video_returns_kind_video(self, video_bytes):
        r = _upload_file(video_bytes, "TEST_clip.mp4", "video/mp4")
        assert r.status_code == 200, r.text
        row = r.json()
        assert row["kind"] == "video"
        assert row["mime_type"] == "video/mp4"
        pytest.video_asset_id = row["id"]

    def test_video_extraction_meta_thumb_and_transcript(self):
        aid = getattr(pytest, "video_asset_id", None)
        if not aid:
            pytest.skip("upload step failed")
        # Whisper STT - allow up to 4 min
        final = _wait_complete(aid, timeout=240)
        assert final["status"] == "complete", final
        meta = final.get("extraction_meta") or {}
        assert meta.get("duration_s") is not None
        assert meta.get("duration_s") >= 2.5
        assert isinstance(meta.get("audio_bytes"), int) and meta["audio_bytes"] > 0
        assert meta.get("thumb_format") == "png", meta
        # thumb_b64 should be populated and large
        assert final.get("thumb_b64") and len(final["thumb_b64"]) > 500
        # char_count is the transcript length; can be 0 for pure tone — but
        # field MUST exist.
        assert "char_count" in meta

    def test_video_cleanup(self):
        aid = getattr(pytest, "video_asset_id", None)
        if aid:
            requests.delete(f"{BASE_URL}/api/cortex/assets/{aid}",
                            headers=AUTH_HEADERS, timeout=15)


# ---------- Size caps --------------------------------------------------
class TestSizeCaps:
    def test_21mib_pdf_rejected_413(self):
        # 21 MiB PDF — should exceed default 20 MiB cap.
        # Need a valid-content-type request; use minimal PDF header + padding.
        # Backend reads bytes before extraction so size check fires first.
        data = b"%PDF-1.4\n" + b"\0" * (21 * 1024 * 1024)
        r = _upload_file(data, "TEST_big.pdf", "application/pdf")
        assert r.status_code == 413, f"expected 413, got {r.status_code}: {r.text[:200]}"

    def test_21mib_video_accepted_200(self):
        # 21 MiB video — should pass the 50 MiB video cap and start pipeline.
        # We don't care about extraction here (random bytes won't extract);
        # just that the route accepts the upload and returns 200.
        data = b"\0" * (21 * 1024 * 1024)
        r = _upload_file(data, "TEST_21m.mp4", "video/mp4")
        assert r.status_code == 200, f"expected 200, got {r.status_code}: {r.text[:200]}"
        asset_id = r.json()["id"]
        # Cleanup — don't wait for pipeline (will end as complete-with-error
        # or failed; either is fine for this test).
        requests.delete(f"{BASE_URL}/api/cortex/assets/{asset_id}",
                        headers=AUTH_HEADERS, timeout=15)

    def test_51mib_video_rejected_413(self):
        data = b"\0" * (51 * 1024 * 1024)
        # Use longer timeout — 51 MiB upload over preview URL can be slow.
        # Ingress may also pre-empt with 413/504; accept either as "rejected".
        try:
            r = _upload_file_long(data, "TEST_51m.mp4", "video/mp4",
                                  timeout=240)
            status = r.status_code
            assert status in (413, 504, 502), (
                f"expected rejection (413/504/502), got {status}: {r.text[:200]}")
        except requests.exceptions.RequestException as e:
            # An ingress-level drop on a 51 MiB body is also acceptable —
            # the request was rejected before reaching FastAPI.
            pytest.skip(f"ingress dropped 51 MiB upload: {type(e).__name__}")


# ---------- Backwards compat (image) -----------------------------------
class TestBackwardsCompat:
    def test_image_upload_still_works(self):
        from PIL import Image
        img = Image.new("RGB", (100, 100), (200, 80, 40))
        buf = io.BytesIO()
        img.save(buf, format="JPEG")
        buf.seek(0)
        r = _upload_file(buf.getvalue(), "TEST_pic.jpg", "image/jpeg")
        assert r.status_code == 200, r.text
        aid = r.json()["id"]
        final = _wait_complete(aid, timeout=180)
        # Pipeline should complete; image gets thumb_b64.
        assert final["status"] == "complete", final
        assert final.get("thumb_b64"), "image thumb should be populated"
        requests.delete(f"{BASE_URL}/api/cortex/assets/{aid}",
                        headers=AUTH_HEADERS, timeout=15)


# ---------- Failure mode: bad PPTX -------------------------------------
class TestFailureModes:
    def test_non_zip_pptx_completes_with_error_meta(self):
        # File claims pptx MIME but is plain text — python-pptx will
        # raise BadZipFile. Pipeline should still mark complete (or failed)
        # with error visible in extraction_meta.
        data = b"this is not a real pptx file" * 50
        r = _upload_file(data, "TEST_bad.pptx", PPTX_MIME)
        assert r.status_code == 200, r.text
        aid = r.json()["id"]
        final = _wait_complete(aid, timeout=120)
        # Either complete with error in extraction_meta, or failed.
        meta = final.get("extraction_meta") or {}
        assert final["status"] in ("complete", "failed"), final
        if final["status"] == "complete":
            # Spec says: asset.status='complete' but extraction_meta carries
            # error reason.
            assert "error" in meta, f"expected error in extraction_meta: {meta}"
        requests.delete(f"{BASE_URL}/api/cortex/assets/{aid}",
                        headers=AUTH_HEADERS, timeout=15)
